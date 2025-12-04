"""
PCA Agent - Automated Reporting Module

Upload report templates (XLSX, CSV, PPTX) and automatically populate them with campaign data.
The system understands template structure and maps data to the correct locations.

Run with:
    streamlit run streamlit_reporting.py --server.port 8504
"""

import os
import sys
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Any
from datetime import datetime
import io

import streamlit as st
import pandas as pd
import numpy as np
from dotenv import load_dotenv

# Add project root to path
current_dir = os.path.dirname(os.path.abspath(__file__))
if current_dir not in sys.path:
    sys.path.insert(0, current_dir)

load_dotenv()

# Page config
st.set_page_config(
    page_title="PCA Agent - Reporting",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded",
)

# Inject CSS
st.markdown("""
<style>
    .main-header {
        font-size: 2.4rem;
        font-weight: 700;
        color: #5b6ef5;
        margin-bottom: 1rem;
    }
    .section-header {
        font-size: 1.6rem;
        font-weight: 600;
        color: #667eea;
        margin-top: 1.5rem;
        margin-bottom: 0.5rem;
    }
    .info-box {
        background-color: rgba(91, 110, 245, 0.1);
        border-left: 4px solid #5b6ef5;
        padding: 1rem;
        margin: 1rem 0;
        border-radius: 4px;
    }
</style>
""", unsafe_allow_html=True)


def init_session_state():
    """Initialize session state variables."""
    defaults = {
        "template_file": None,
        "template_type": None,
        "template_structure": None,
        "data_file": None,
        "data_df": None,
        "mapping_config": {},
        "generated_report": None,
    }
    for key, value in defaults.items():
        if key not in st.session_state:
            st.session_state[key] = value


def render_sidebar() -> str:
    """Render sidebar navigation."""
    with st.sidebar:
        st.markdown("## 📊 Reporting Module")
        st.markdown("*Automated Report Generation*")
        
        st.info("""
        Upload your source data and report template to automatically generate 
        populated reports in the same format.
        """)
        
        st.divider()
        
        page = st.radio(
            "Navigation",
            options=["Upload Data & Template", "Data Mapping", "Generate Report", "Settings"],
            index=0,
        )
        
        st.divider()
        
        # Status indicators
        st.subheader("Status")
        
        if st.session_state.template_file:
            st.success(f"✅ Template: {st.session_state.template_type}")
        else:
            st.warning("⚠️ No template uploaded")
        
        if st.session_state.data_df is not None:
            st.success(f"✅ Data: {len(st.session_state.data_df)} rows")
        else:
            st.warning("⚠️ No data loaded")
        
        if st.session_state.generated_report:
            st.success("✅ Report generated")
        
        st.divider()
        
        if st.button("🔄 Reset All"):
            for key in ["template_file", "template_type", "template_structure", 
                       "data_file", "data_df", "mapping_config", "generated_report"]:
                st.session_state[key] = None if key != "mapping_config" else {}
            st.rerun()
    
    return page.lower().replace(" ", "_")


def render_upload_page():
    """Render combined data and template upload page."""
    st.markdown('<div class="main-header">📤 Upload Data & Template</div>', unsafe_allow_html=True)
    
    st.markdown("""
    ### Step 1: Upload Source Data
    Upload your campaign performance data that will populate the report template.
    """)
    
    # Data upload section
    col1, col2 = st.columns([2, 1])
    
    with col1:
        data_file = st.file_uploader(
            "📊 Upload Source Data",
            type=["csv", "xlsx", "xls"],
            help="Upload your campaign performance data (CSV or Excel)",
            key="data_uploader"
        )
    
    with col2:
        st.markdown("#### Supported Formats")
        st.markdown("- CSV (.csv)")
        st.markdown("- Excel (.xlsx, .xls)")
    
    if data_file is not None:
        try:
            # Load data
            if data_file.name.endswith('.csv'):
                df = pd.read_csv(data_file)
            else:
                df = pd.read_excel(data_file)
            
            st.session_state.data_df = df
            st.session_state.data_file = data_file.name
            
            st.success(f"✅ Data loaded: {len(df)} rows, {len(df.columns)} columns")
            
            # Data preview
            with st.expander("📋 Data Preview", expanded=True):
                col1, col2, col3, col4 = st.columns(4)
                with col1:
                    st.metric("Rows", len(df))
                with col2:
                    st.metric("Columns", len(df.columns))
                with col3:
                    numeric_cols = len(df.select_dtypes(include=['number']).columns)
                    st.metric("Numeric Cols", numeric_cols)
                with col4:
                    text_cols = len(df.select_dtypes(include=['object']).columns)
                    st.metric("Text Cols", text_cols)
                
                st.dataframe(df.head(10), use_container_width=True)
                
                # Column info
                with st.expander("ℹ️ Column Information"):
                    col_info = pd.DataFrame({
                        'Column': df.columns,
                        'Type': df.dtypes.astype(str),
                        'Non-Null': df.count().values,
                        'Null': df.isnull().sum().values,
                        'Unique': [df[col].nunique() for col in df.columns]
                    })
                    st.dataframe(col_info, use_container_width=True)
        
        except Exception as e:
            st.error(f"❌ Error loading data: {str(e)}")
    
    st.divider()
    
    # Template upload section
    st.markdown("""
    ### Step 2: Upload Report Template
    
    Upload your report template with placeholders. Supported formats:
    - **Excel (.xlsx, .xls)** - Spreadsheet templates with formulas and formatting
    - **CSV (.csv)** - Simple tabular templates
    - **PowerPoint (.pptx)** - Presentation templates with charts and tables
    
    The system will analyze the template structure and identify data placeholders.
    """)
    
    col1, col2 = st.columns([2, 1])
    
    with col1:
        template_file = st.file_uploader(
            "📄 Upload Report Template",
            type=["xlsx", "xls", "csv", "pptx"],
            help="Upload your report template in Excel, CSV, or PowerPoint format",
            key="template_uploader"
        )
    
    with col2:
        st.markdown("#### Placeholder Formats")
        st.markdown("- `{{field}}`")
        st.markdown("- `{field}`")
        st.markdown("- `[field]`")
        st.markdown("- `<field>`")
    
    if template_file is not None:
        file_ext = template_file.name.split('.')[-1].lower()
        st.session_state.template_file = template_file
        st.session_state.template_type = file_ext
        
        st.success(f"✅ Template uploaded: {template_file.name}")
        
        # Analyze template
        with st.spinner("🔍 Analyzing template structure..."):
            structure = analyze_template(template_file, file_ext)
            st.session_state.template_structure = structure
        
        # Display template analysis
        st.markdown('<div class="section-header">Template Analysis</div>', unsafe_allow_html=True)
        
        col1, col2, col3 = st.columns(3)
        
        with col1:
            st.metric("Format", file_ext.upper())
        
        with col2:
            if file_ext in ['xlsx', 'xls', 'csv']:
                st.metric("Sheets/Tables", structure.get('sheet_count', 1))
        
        with col3:
            st.metric("Data Fields", structure.get('field_count', 0))
        
        # Show detected placeholders
        if structure.get('placeholders'):
            with st.expander("🔍 Detected Placeholders", expanded=True):
                st.markdown("These fields were detected in your template:")
                for placeholder in structure['placeholders']:
                    st.markdown(f"- `{placeholder['name']}` ({placeholder['type']}) - {placeholder['location']}")
        
        # Show preview
        if file_ext in ['xlsx', 'xls', 'csv']:
            with st.expander("📋 Template Preview", expanded=False):
                preview_df = load_template_preview(template_file, file_ext)
                if preview_df is not None:
                    st.dataframe(preview_df.head(10), use_container_width=True)


def analyze_template(file, file_type: str) -> Dict[str, Any]:
    """
    Analyze template structure and identify placeholders.
    
    Args:
        file: Uploaded file object
        file_type: File extension (xlsx, csv, pptx)
    
    Returns:
        Dictionary with template structure information
    """
    structure = {
        'file_type': file_type,
        'placeholders': [],
        'field_count': 0,
        'sheet_count': 0
    }
    
    try:
        if file_type in ['xlsx', 'xls']:
            import openpyxl
            
            # Load workbook
            wb = openpyxl.load_workbook(file, data_only=False)
            structure['sheet_count'] = len(wb.sheetnames)
            
            placeholders = []
            
            # Scan each sheet for placeholders
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value and isinstance(cell.value, str):
                            # Detect placeholder patterns: {{field}}, {field}, [field], <field>
                            if any(pattern in str(cell.value) for pattern in ['{{', '{', '[', '<']):
                                placeholders.append({
                                    'name': cell.value,
                                    'type': 'text',
                                    'location': f"{sheet_name}!{cell.coordinate}",
                                    'sheet': sheet_name,
                                    'cell': cell.coordinate
                                })
            
            structure['placeholders'] = placeholders
            structure['field_count'] = len(placeholders)
            
        elif file_type == 'csv':
            df = pd.read_csv(file)
            structure['sheet_count'] = 1
            
            # Check for placeholder columns
            placeholders = []
            for col in df.columns:
                if any(pattern in str(col) for pattern in ['{{', '{', '[', '<']):
                    placeholders.append({
                        'name': col,
                        'type': 'column',
                        'location': f"Column: {col}",
                        'sheet': 'main',
                        'cell': col
                    })
            
            structure['placeholders'] = placeholders
            structure['field_count'] = len(placeholders)
            
        elif file_type == 'pptx':
            from pptx import Presentation
            
            prs = Presentation(file)
            structure['sheet_count'] = len(prs.slides)
            
            placeholders = []
            
            for slide_idx, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        if any(pattern in shape.text for pattern in ['{{', '{', '[', '<']):
                            placeholders.append({
                                'name': shape.text,
                                'type': 'text',
                                'location': f"Slide {slide_idx + 1}",
                                'sheet': f"slide_{slide_idx}",
                                'cell': shape.name
                            })
            
            structure['placeholders'] = placeholders
            structure['field_count'] = len(placeholders)
    
    except Exception as e:
        st.error(f"Error analyzing template: {str(e)}")
    
    return structure


def load_template_preview(file, file_type: str) -> Optional[pd.DataFrame]:
    """Load template preview as DataFrame."""
    try:
        if file_type in ['xlsx', 'xls']:
            return pd.read_excel(file, nrows=10)
        elif file_type == 'csv':
            return pd.read_csv(file, nrows=10)
    except Exception as e:
        st.warning(f"Could not load preview: {str(e)}")
    return None


def render_data_mapping_page():
    """Render data mapping configuration page."""
    st.markdown('<div class="main-header">🔗 Data Mapping</div>', unsafe_allow_html=True)
    
    if not st.session_state.template_file:
        st.warning("⚠️ Please upload a template first in the 'Upload Data & Template' page")
        return
    
    if st.session_state.data_df is None:
        st.warning("⚠️ Please upload source data first in the 'Upload Data & Template' page")
        return
    
    st.markdown("""
    ### Configure Field Mapping
    
    Map template placeholders to your data columns. The system has auto-suggested matches based on field names.
    """)
    
    st.divider()
    
    df = st.session_state.data_df
    
    # Show data summary
    col1, col2, col3 = st.columns(3)
    with col1:
        st.metric("Data Source", st.session_state.data_file)
    with col2:
        st.metric("Rows", len(df))
    with col3:
        st.metric("Columns", len(df.columns))
    
    st.divider()
    
    # Mapping configuration
    st.markdown('<div class="section-header">Field Mappings</div>', unsafe_allow_html=True)
    
    if st.session_state.template_structure and st.session_state.template_structure.get('placeholders'):
        st.markdown("Map template placeholders to data columns:")
        
        mapping_config = {}
        
        for placeholder in st.session_state.template_structure['placeholders']:
            col1, col2, col3 = st.columns([2, 2, 1])
            
            with col1:
                st.text_input(
                    "Template Field",
                    value=placeholder['name'],
                    disabled=True,
                    key=f"template_{placeholder['location']}"
                )
            
            with col2:
                # Auto-suggest matching column
                suggested_col = find_matching_column(placeholder['name'], df.columns)
                
                data_column = st.selectbox(
                    "Data Column",
                    options=['None'] + list(df.columns),
                    index=list(df.columns).index(suggested_col) + 1 if suggested_col else 0,
                    key=f"mapping_{placeholder['location']}"
                )
                
                if data_column != 'None':
                    mapping_config[placeholder['name']] = {
                        'data_column': data_column,
                        'location': placeholder['location'],
                        'type': placeholder['type']
                    }
            
            with col3:
                st.caption(placeholder['location'])
        
        st.session_state.mapping_config = mapping_config
        
        if st.button("💾 Save Mapping", type="primary"):
            st.success("✅ Mapping configuration saved!")
    else:
        st.info("No placeholders detected in template. The entire data will be inserted.")


def find_matching_column(placeholder: str, columns: List[str]) -> Optional[str]:
    """Find best matching column for a placeholder."""
    # Clean placeholder
    clean_placeholder = placeholder.strip('{}[]<>').lower().replace('_', ' ').replace('-', ' ')
    
    # Try exact match
    for col in columns:
        if col.lower() == clean_placeholder:
            return col
    
    # Try partial match
    for col in columns:
        if clean_placeholder in col.lower() or col.lower() in clean_placeholder:
            return col
    
    return None


def render_generate_report_page():
    """Render report generation page."""
    st.markdown('<div class="main-header">📊 Generate Report</div>', unsafe_allow_html=True)
    
    if not st.session_state.template_file:
        st.warning("⚠️ Please upload a template first")
        return
    
    if st.session_state.data_df is None:
        st.warning("⚠️ Please upload data and configure mapping")
        return
    
    st.markdown("""
    ### Generate Your Report
    
    Review the configuration and generate your populated report.
    """)
    
    st.divider()
    
    # Configuration summary
    col1, col2, col3 = st.columns(3)
    
    with col1:
        st.metric("Template", st.session_state.template_type.upper())
    
    with col2:
        st.metric("Data Rows", len(st.session_state.data_df))
    
    with col3:
        st.metric("Mapped Fields", len(st.session_state.mapping_config))
    
    st.divider()
    
    # Generation options
    st.markdown('<div class="section-header">Generation Options</div>', unsafe_allow_html=True)
    
    col1, col2 = st.columns(2)
    
    with col1:
        report_name = st.text_input(
            "Report Name",
            value=f"Report_{datetime.now().strftime('%Y%m%d_%H%M%S')}",
            help="Name for the generated report"
        )
    
    with col2:
        aggregation_method = st.selectbox(
            "Data Aggregation",
            options=["Sum", "Average", "Latest", "All Rows"],
            help="How to aggregate data if multiple rows match"
        )
    
    include_charts = st.checkbox("Include Charts", value=True, help="Generate charts from data")
    include_summary = st.checkbox("Include Summary Sheet", value=True, help="Add executive summary")
    
    st.divider()
    
    # Generate button
    if st.button("🚀 Generate Report", type="primary", use_container_width=True):
        with st.spinner("📊 Generating report..."):
            try:
                generated_file = generate_report(
                    template_file=st.session_state.template_file,
                    template_type=st.session_state.template_type,
                    data_df=st.session_state.data_df,
                    mapping_config=st.session_state.mapping_config,
                    report_name=report_name,
                    aggregation=aggregation_method.lower(),
                    include_charts=include_charts,
                    include_summary=include_summary
                )
                
                st.session_state.generated_report = generated_file
                st.success("✅ Report generated successfully!")
                
            except Exception as e:
                st.error(f"❌ Error generating report: {str(e)}")
    
    # Download generated report
    if st.session_state.generated_report:
        st.divider()
        st.markdown('<div class="section-header">Download Report</div>', unsafe_allow_html=True)
        
        file_ext = st.session_state.template_type
        mime_types = {
            'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            'xls': 'application/vnd.ms-excel',
            'csv': 'text/csv',
            'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        }
        
        st.download_button(
            label=f"📥 Download {report_name}.{file_ext}",
            data=st.session_state.generated_report,
            file_name=f"{report_name}.{file_ext}",
            mime=mime_types.get(file_ext, 'application/octet-stream'),
            use_container_width=True
        )


def generate_report(
    template_file,
    template_type: str,
    data_df: pd.DataFrame,
    mapping_config: Dict,
    report_name: str,
    aggregation: str = 'sum',
    include_charts: bool = True,
    include_summary: bool = True
) -> bytes:
    """
    Generate populated report from template and data.
    
    Args:
        template_file: Template file object
        template_type: File extension
        data_df: Campaign data
        mapping_config: Field mapping configuration
        report_name: Output report name
        aggregation: Data aggregation method
        include_charts: Whether to include charts
        include_summary: Whether to include summary
    
    Returns:
        Generated report as bytes
    """
    if template_type in ['xlsx', 'xls']:
        return generate_excel_report(
            template_file, data_df, mapping_config, aggregation, include_charts, include_summary
        )
    elif template_type == 'csv':
        return generate_csv_report(data_df, mapping_config)
    elif template_type == 'pptx':
        return generate_pptx_report(
            template_file, data_df, mapping_config, include_charts
        )
    else:
        raise ValueError(f"Unsupported template type: {template_type}")


def generate_excel_report(
    template_file,
    data_df: pd.DataFrame,
    mapping_config: Dict,
    aggregation: str,
    include_charts: bool,
    include_summary: bool
) -> bytes:
    """
    Generate Excel report by replacing placeholders with actual data.
    
    Process:
    1. Load template workbook
    2. Find all placeholders (e.g., {{Total_Spend}})
    3. Replace with aggregated data from corresponding column
    4. Preserve all formatting, formulas, and charts
    5. Add raw data sheet if needed
    """
    import openpyxl
    from openpyxl.utils.dataframe import dataframe_to_rows
    from openpyxl.styles import Font, PatternFill
    
    # Load template (preserving formatting)
    wb = openpyxl.load_workbook(template_file)
    
    # Replace placeholders with actual data
    for placeholder_name, config in mapping_config.items():
        data_column = config['data_column']
        location = config['location']
        
        # Parse location (e.g., "Sheet1!A1")
        if '!' in location:
            sheet_name, cell_ref = location.split('!')
            sheet = wb[sheet_name]
            
            # Calculate aggregated value based on method
            try:
                if aggregation == 'sum':
                    value = data_df[data_column].sum()
                elif aggregation == 'average':
                    value = data_df[data_column].mean()
                elif aggregation == 'latest':
                    value = data_df[data_column].iloc[-1] if len(data_df) > 0 else 0
                elif aggregation == 'all rows':
                    # For "all rows", insert comma-separated values
                    value = ', '.join(map(str, data_df[data_column].tolist()))
                else:
                    value = data_df[data_column].sum()
                
                # Format numeric values
                if isinstance(value, (int, float)):
                    value = round(value, 2)
                
                # Replace placeholder in cell
                cell = sheet[cell_ref]
                
                # If cell contains placeholder pattern, replace it
                if cell.value and isinstance(cell.value, str):
                    cell.value = cell.value.replace(placeholder_name, str(value))
                else:
                    cell.value = value
                    
            except Exception as e:
                print(f"Error processing {placeholder_name}: {e}")
                sheet[cell_ref] = f"Error: {str(e)}"
    
    # Add raw data sheet if no placeholders or if requested
    if not mapping_config or include_summary:
        # Create or get Data sheet
        if 'Campaign_Data' not in wb.sheetnames:
            data_sheet = wb.create_sheet('Campaign_Data')
        else:
            data_sheet = wb['Campaign_Data']
        
        # Clear existing data
        data_sheet.delete_rows(1, data_sheet.max_row)
        
        # Add headers with formatting
        for c_idx, col_name in enumerate(data_df.columns, 1):
            cell = data_sheet.cell(row=1, column=c_idx, value=col_name)
            cell.font = Font(bold=True, color="FFFFFF")
            cell.fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
        
        # Add data rows
        for r_idx, row in enumerate(dataframe_to_rows(data_df, index=False, header=False), 2):
            for c_idx, value in enumerate(row, 1):
                data_sheet.cell(row=r_idx, column=c_idx, value=value)
        
        # Auto-adjust column widths
        for column in data_sheet.columns:
            max_length = 0
            column_letter = column[0].column_letter
            for cell in column:
                try:
                    if len(str(cell.value)) > max_length:
                        max_length = len(str(cell.value))
                except:
                    pass
            adjusted_width = min(max_length + 2, 50)
            data_sheet.column_dimensions[column_letter].width = adjusted_width
    
    # Save to bytes
    output = io.BytesIO()
    wb.save(output)
    output.seek(0)
    return output.getvalue()


def generate_csv_report(data_df: pd.DataFrame, mapping_config: Dict) -> bytes:
    """Generate CSV report."""
    output = io.StringIO()
    data_df.to_csv(output, index=False)
    return output.getvalue().encode('utf-8')


def generate_pptx_report(
    template_file,
    data_df: pd.DataFrame,
    mapping_config: Dict,
    include_charts: bool
) -> bytes:
    """
    Generate PowerPoint report by replacing placeholders with actual data.
    
    Process:
    1. Load template presentation
    2. Scan all slides for text placeholders
    3. Replace placeholders with aggregated data
    4. Preserve all formatting, layouts, and existing charts
    5. Optionally add data summary slide
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    
    # Load template (preserving all formatting and layouts)
    prs = Presentation(template_file)
    
    # Replace text placeholders in all slides
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            # Check text boxes and text frames
            if hasattr(shape, "text") and shape.text:
                original_text = shape.text
                updated_text = original_text
                
                # Replace all mapped placeholders
                for placeholder_name, config in mapping_config.items():
                    if placeholder_name in updated_text:
                        data_column = config['data_column']
                        
                        try:
                            # Calculate aggregated value
                            value = data_df[data_column].sum()
                            
                            # Format numeric values
                            if isinstance(value, (int, float)):
                                if value >= 1000000:
                                    formatted_value = f"${value/1000000:.2f}M"
                                elif value >= 1000:
                                    formatted_value = f"${value/1000:.1f}K"
                                else:
                                    formatted_value = f"${value:,.2f}"
                            else:
                                formatted_value = str(value)
                            
                            # Replace placeholder
                            updated_text = updated_text.replace(placeholder_name, formatted_value)
                        except Exception as e:
                            print(f"Error processing {placeholder_name} on slide {slide_idx + 1}: {e}")
                            updated_text = updated_text.replace(placeholder_name, "N/A")
                
                # Update text if changed
                if updated_text != original_text:
                    if hasattr(shape, "text_frame"):
                        shape.text_frame.text = updated_text
                    else:
                        shape.text = updated_text
            
            # Check tables for placeholders
            if hasattr(shape, "table"):
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text:
                            original_text = cell.text
                            updated_text = original_text
                            
                            for placeholder_name, config in mapping_config.items():
                                if placeholder_name in updated_text:
                                    data_column = config['data_column']
                                    try:
                                        value = data_df[data_column].sum()
                                        if isinstance(value, (int, float)):
                                            value = round(value, 2)
                                        updated_text = updated_text.replace(placeholder_name, str(value))
                                    except:
                                        updated_text = updated_text.replace(placeholder_name, "N/A")
                            
                            if updated_text != original_text:
                                cell.text = updated_text
    
    # Save to bytes
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()


def render_settings_page():
    """Render settings page."""
    st.markdown('<div class="main-header">⚙️ Settings</div>', unsafe_allow_html=True)
    
    st.markdown("""
    ### Report Generation Settings
    
    Configure default behaviors for report generation.
    """)
    
    st.divider()
    
    st.subheader("Default Options")
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.selectbox("Default Aggregation", ["Sum", "Average", "Latest", "All Rows"])
        st.checkbox("Auto-detect placeholders", value=True)
    
    with col2:
        st.selectbox("Date Format", ["YYYY-MM-DD", "DD/MM/YYYY", "MM/DD/YYYY"])
        st.checkbox("Include metadata", value=True)
    
    st.divider()
    
    st.subheader("Template Library")
    st.info("Coming soon: Save and reuse template configurations")


def main():
    """Main application entry point."""
    init_session_state()
    page = render_sidebar()
    
    if page == "upload_data_&_template":
        render_upload_page()
    elif page == "data_mapping":
        render_data_mapping_page()
    elif page == "generate_report":
        render_generate_report_page()
    elif page == "settings":
        render_settings_page()
    
    st.divider()
    st.caption("PCA Agent Reporting Module · Automated Report Generation")


if __name__ == "__main__":
    main()
